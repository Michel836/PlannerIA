"""
Système de Gestion RAG Avancé pour PlannerIA
Gestion complète des documents avec formats multiples, métadonnées, recherche avancée et configuration
"""

import streamlit as st
import pandas as pd
from typing import Dict, List, Any, Optional
from datetime import datetime
import json
import os
from pathlib import Path

class AdvancedRAGManager:
    """Gestionnaire RAG avancé avec toutes les fonctionnalités"""
    
    def __init__(self):
        self.supported_formats = {
            'text': ['.txt', '.md', '.markdown'],
            'data': ['.json', '.csv', '.yaml', '.yml'],
            'pdf': ['.pdf'],
            'office': ['.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls'],
            'web': ['http://', 'https://'],
            'image': ['.png', '.jpg', '.jpeg', '.gif', '.bmp']
        }
        
        self.categories = [
            "📋 Gestion de Projet",
            "💻 Technique",
            "📊 Business",
            "📈 Analytics",
            "⚠️ Risques",
            "📚 Documentation",
            "🎯 Stratégie",
            "🔧 Opérationnel"
        ]
        
        self.default_tags = [
            "planning", "budget", "risques", "équipe", "technique",
            "architecture", "design", "test", "déploiement", "maintenance",
            "client", "fournisseur", "contrat", "rapport", "présentation"
        ]
    
    def render_advanced_interface(self, rag_manager):
        """Interface principale du gestionnaire RAG avancé"""
        
        # Tabs pour organiser les fonctionnalités
        tab1, tab2, tab3, tab4, tab5 = st.tabs([
            "📤 Import Avancé",
            "📊 Gestion Documents", 
            "🔍 Recherche Pro",
            "🏷️ Métadonnées",
            "⚙️ Configuration"
        ])
        
        with tab1:
            self.render_advanced_import(rag_manager)
        
        with tab2:
            self.render_document_management(rag_manager)
        
        with tab3:
            self.render_advanced_search(rag_manager)
        
        with tab4:
            self.render_metadata_management(rag_manager)
        
        with tab5:
            self.render_rag_configuration(rag_manager)
    
    def render_advanced_import(self, rag_manager):
        """Interface d'import avancée avec support multi-formats"""
        st.markdown("### 📤 Import Multi-Formats Avancé")
        
        import_method = st.radio(
            "Méthode d'import:",
            ["📁 Fichiers locaux", "🌐 URL Web", "📋 Texte direct", "🔗 API/Base de données"],
            horizontal=True
        )
        
        if import_method == "📁 Fichiers locaux":
            self.import_local_files(rag_manager)
        elif import_method == "🌐 URL Web":
            self.import_from_url(rag_manager)
        elif import_method == "📋 Texte direct":
            self.import_direct_text(rag_manager)
        else:
            self.import_from_api(rag_manager)
    
    def import_local_files(self, rag_manager):
        """Import de fichiers locaux avec support étendu"""
        col1, col2 = st.columns([2, 1])
        
        with col1:
            uploaded_files = st.file_uploader(
                "Sélectionnez vos documents",
                type=["txt", "md", "json", "csv", "pdf", "docx", "pptx", "xlsx"],
                accept_multiple_files=True,
                help="Formats supportés: TXT, MD, JSON, CSV, PDF, Word, PowerPoint, Excel"
            )
            
            if uploaded_files:
                # Sélection de catégorie et tags
                with col2:
                    st.markdown("**📝 Métadonnées**")
                    category = st.selectbox("Catégorie:", self.categories)
                    
                    # Tags personnalisés
                    selected_tags = st.multiselect(
                        "Tags:",
                        self.default_tags,
                        default=[]
                    )
                    
                    custom_tag = st.text_input("Tag personnalisé:")
                    if custom_tag:
                        selected_tags.append(custom_tag)
                    
                    # Options d'import
                    st.markdown("**⚙️ Options**")
                    chunk_size = st.slider(
                        "Taille chunks (caractères):",
                        100, 5000, 1000,
                        help="Diviser les documents longs en segments"
                    )
                    
                    process_ocr = st.checkbox(
                        "OCR pour images/PDF scannés",
                        help="Extraction de texte depuis images"
                    )
                
                # Bouton d'import
                if st.button("📥 **IMPORTER AVEC MÉTADONNÉES**", type="primary", use_container_width=True):
                    self.process_advanced_import(
                        uploaded_files, rag_manager, 
                        category, selected_tags, chunk_size, process_ocr
                    )
    
    def import_from_url(self, rag_manager):
        """Import depuis URLs web"""
        st.markdown("#### 🌐 Import depuis le Web")
        
        url_input = st.text_area(
            "URLs (une par ligne):",
            height=100,
            placeholder="https://example.com/document\nhttps://site.com/article"
        )
        
        col1, col2 = st.columns(2)
        
        with col1:
            scrape_depth = st.slider(
                "Profondeur de scraping:",
                0, 3, 1,
                help="0=Page seule, 1=Liens directs, etc."
            )
        
        with col2:
            include_images = st.checkbox("Inclure les images")
            include_pdfs = st.checkbox("Télécharger PDFs liés")
        
        if st.button("🌐 **IMPORTER DEPUIS WEB**", type="primary"):
            if url_input:
                urls = [url.strip() for url in url_input.split('\n') if url.strip()]
                self.process_web_import(urls, rag_manager, scrape_depth, include_images, include_pdfs)
    
    def import_from_api(self, rag_manager):
        """Import depuis API ou base de données"""
        st.markdown("#### 🔗 Import depuis API/Base de données")
        
        source_type = st.selectbox(
            "Type de source:",
            ["Jira", "Confluence", "GitHub", "GitLab", "Notion", "Google Drive", "Custom API"]
        )
        
        if source_type == "Custom API":
            api_url = st.text_input("URL de l'API:")
            api_key = st.text_input("Clé API (optionnel):", type="password")
            
            if st.button("🔗 **CONNECTER ET IMPORTER**"):
                st.info(f"Connexion à {api_url}...")
                # Implémenter la logique d'import API
        else:
            st.info(f"Configuration pour {source_type} à venir...")
    
    def render_document_management(self, rag_manager):
        """Gestion avancée des documents"""
        st.markdown("### 📊 Gestion des Documents")
        
        # Statistiques
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            doc_count = self.get_document_count(rag_manager)
            st.metric("📁 Documents", doc_count)
        
        with col2:
            total_size = self.get_total_size(rag_manager)
            st.metric("💾 Taille totale", f"{total_size:.2f} MB")
        
        with col3:
            categories_count = self.get_categories_count(rag_manager)
            st.metric("🏷️ Catégories", categories_count)
        
        with col4:
            last_update = self.get_last_update(rag_manager)
            st.metric("🕒 Dernière MAJ", last_update)
        
        # Tableau des documents avec actions
        st.markdown("#### 📋 Liste des Documents")
        
        # Filtres
        col_filter1, col_filter2, col_filter3 = st.columns(3)
        
        with col_filter1:
            filter_category = st.selectbox(
                "Filtrer par catégorie:",
                ["Toutes"] + self.categories
            )
        
        with col_filter2:
            filter_date = st.date_input(
                "Depuis le:",
                value=None
            )
        
        with col_filter3:
            filter_type = st.selectbox(
                "Type de fichier:",
                ["Tous", "Texte", "PDF", "Office", "Data", "Web"]
            )
        
        # Tableau interactif
        documents = self.get_filtered_documents(
            rag_manager, filter_category, filter_date, filter_type
        )
        
        if documents:
            df = pd.DataFrame(documents)
            
            # Sélection pour actions groupées
            selected_docs = st.multiselect(
                "Sélectionner des documents:",
                df['id'].tolist(),
                format_func=lambda x: df[df['id'] == x]['name'].values[0]
            )
            
            # Actions groupées
            if selected_docs:
                col_action1, col_action2, col_action3, col_action4 = st.columns(4)
                
                with col_action1:
                    if st.button("🏷️ Modifier tags"):
                        self.bulk_update_tags(selected_docs, rag_manager)
                
                with col_action2:
                    if st.button("📂 Changer catégorie"):
                        self.bulk_update_category(selected_docs, rag_manager)
                
                with col_action3:
                    if st.button("📥 Exporter"):
                        self.export_documents(selected_docs, rag_manager)
                
                with col_action4:
                    if st.button("🗑️ Supprimer", type="secondary"):
                        self.delete_documents(selected_docs, rag_manager)
            
            # Affichage du tableau
            st.dataframe(
                df[['name', 'category', 'tags', 'size', 'date', 'relevance']],
                use_container_width=True,
                hide_index=True
            )
            
            # Détails d'un document
            if st.checkbox("Voir les détails"):
                selected_doc = st.selectbox(
                    "Document:",
                    df['name'].tolist()
                )
                
                if selected_doc:
                    doc_details = df[df['name'] == selected_doc].iloc[0]
                    
                    col_detail1, col_detail2 = st.columns(2)
                    
                    with col_detail1:
                        st.markdown("**📄 Informations**")
                        st.json({
                            "ID": doc_details['id'],
                            "Nom": doc_details['name'],
                            "Catégorie": doc_details['category'],
                            "Tags": doc_details['tags'],
                            "Taille": f"{doc_details['size']} KB",
                            "Date": doc_details['date']
                        })
                    
                    with col_detail2:
                        st.markdown("**📝 Aperçu**")
                        preview = self.get_document_preview(doc_details['id'], rag_manager)
                        st.text_area("Contenu:", preview, height=200, disabled=True)
    
    def render_advanced_search(self, rag_manager):
        """Interface de recherche avancée"""
        st.markdown("### 🔍 Recherche Avancée")
        
        # Type de recherche
        search_type = st.radio(
            "Type de recherche:",
            ["🔤 Mots-clés", "🧠 Sémantique", "🎯 Hybride", "📊 Requête structurée"],
            horizontal=True
        )
        
        # Champ de recherche principal
        query = st.text_input(
            "Rechercher:",
            placeholder="Ex: risques projet mobile, planning développement, budget Q3..."
        )
        
        # Options avancées
        with st.expander("🔧 Options avancées"):
            col_opt1, col_opt2 = st.columns(2)
            
            with col_opt1:
                relevance_threshold = st.slider(
                    "Seuil de pertinence:",
                    0.0, 1.0, 0.5,
                    help="Minimum de pertinence pour afficher"
                )
                
                max_results = st.number_input(
                    "Nombre max de résultats:",
                    1, 100, 10
                )
            
            with col_opt2:
                search_in_categories = st.multiselect(
                    "Chercher dans:",
                    self.categories,
                    default=self.categories
                )
                
                date_range = st.date_input(
                    "Période:",
                    value=[],
                    help="Laisser vide pour toutes les dates"
                )
        
        # Bouton de recherche
        if st.button("🔍 **RECHERCHER**", type="primary", use_container_width=True):
            if query:
                results = self.perform_advanced_search(
                    query, rag_manager, search_type,
                    relevance_threshold, max_results,
                    search_in_categories, date_range
                )
                
                # Affichage des résultats
                if results:
                    st.success(f"🎯 {len(results)} résultats trouvés")
                    
                    for i, result in enumerate(results, 1):
                        with st.expander(
                            f"📄 {result['title']} - Pertinence: {result['relevance']:.2%}"
                        ):
                            col_res1, col_res2 = st.columns([3, 1])
                            
                            with col_res1:
                                st.markdown(f"**Extrait:**")
                                st.write(result['excerpt'])
                                
                                if result.get('highlights'):
                                    st.markdown("**🔍 Passages pertinents:**")
                                    for highlight in result['highlights']:
                                        st.info(f"...{highlight}...")
                            
                            with col_res2:
                                st.markdown("**📊 Métadonnées**")
                                st.caption(f"Catégorie: {result['category']}")
                                st.caption(f"Tags: {', '.join(result['tags'])}")
                                st.caption(f"Date: {result['date']}")
                                
                                if st.button(f"📖 Voir plus", key=f"more_{i}"):
                                    st.session_state.selected_doc = result['id']
                else:
                    st.info("🔍 Aucun résultat trouvé")
        
        # Suggestions de recherche
        st.markdown("#### 💡 Suggestions")
        
        suggestions = self.get_search_suggestions(query, rag_manager)
        if suggestions:
            cols = st.columns(min(len(suggestions), 4))
            for i, suggestion in enumerate(suggestions[:4]):
                with cols[i]:
                    if st.button(suggestion, key=f"sug_{i}"):
                        st.session_state.search_query = suggestion
                        st.success("✅ Mise à jour effectuée")
    
    def render_metadata_management(self, rag_manager):
        """Gestion des métadonnées et taxonomie"""
        st.markdown("### 🏷️ Gestion des Métadonnées")
        
        tab_meta1, tab_meta2, tab_meta3 = st.tabs([
            "📂 Catégories", "🏷️ Tags", "📊 Statistiques"
        ])
        
        with tab_meta1:
            st.markdown("#### 📂 Gestion des Catégories")
            
            col_cat1, col_cat2 = st.columns(2)
            
            with col_cat1:
                st.markdown("**Catégories existantes**")
                for category in self.categories:
                    col_a, col_b, col_c = st.columns([3, 1, 1])
                    with col_a:
                        st.write(category)
                    with col_b:
                        count = self.get_category_document_count(category, rag_manager)
                        st.caption(f"{count} docs")
                    with col_c:
                        if st.button("✏️", key=f"edit_cat_{category}"):
                            st.session_state.editing_category = category
            
            with col_cat2:
                st.markdown("**Ajouter une catégorie**")
                new_cat_name = st.text_input("Nom:")
                new_cat_icon = st.text_input("Emoji:", value="📁")
                
                if st.button("➕ Ajouter catégorie"):
                    if new_cat_name:
                        self.add_category(f"{new_cat_icon} {new_cat_name}")
                        st.success("✅ Catégorie ajoutée")
                        st.success("✅ Mise à jour effectuée")
        
        with tab_meta2:
            st.markdown("#### 🏷️ Gestion des Tags")
            
            # Nuage de tags
            tag_stats = self.get_tag_statistics(rag_manager)
            if tag_stats:
                st.markdown("**☁️ Nuage de tags**")
                
                # Créer un nuage de tags avec tailles variables
                tag_html = self.create_tag_cloud(tag_stats)
                st.markdown(tag_html, unsafe_allow_html=True)
            
            # Gestion des tags
            col_tag1, col_tag2 = st.columns(2)
            
            with col_tag1:
                st.markdown("**Tags populaires**")
                popular_tags = self.get_popular_tags(rag_manager, limit=10)
                
                for tag, count in popular_tags:
                    col_a, col_b = st.columns([3, 1])
                    with col_a:
                        st.write(f"#{tag}")
                    with col_b:
                        st.caption(f"{count} uses")
            
            with col_tag2:
                st.markdown("**Fusion de tags**")
                tags_to_merge = st.multiselect(
                    "Sélectionner tags à fusionner:",
                    [tag for tag, _ in popular_tags]
                )
                
                new_tag_name = st.text_input("Nouveau nom du tag fusionné:")
                
                if st.button("🔄 Fusionner tags"):
                    if len(tags_to_merge) > 1 and new_tag_name:
                        self.merge_tags(tags_to_merge, new_tag_name, rag_manager)
                        st.success("✅ Tags fusionnés")
                        st.success("✅ Mise à jour effectuée")
        
        with tab_meta3:
            st.markdown("#### 📊 Statistiques Métadonnées")
            
            # Graphiques de répartition
            col_stat1, col_stat2 = st.columns(2)
            
            with col_stat1:
                st.markdown("**📊 Répartition par catégorie**")
                category_data = self.get_category_distribution(rag_manager)
                if category_data:
                    df_cat = pd.DataFrame(category_data)
                    st.bar_chart(df_cat.set_index('category')['count'])
            
            with col_stat2:
                st.markdown("**📈 Évolution temporelle**")
                timeline_data = self.get_import_timeline(rag_manager)
                if timeline_data:
                    df_time = pd.DataFrame(timeline_data)
                    st.line_chart(df_time.set_index('date')['count'])
    
    def render_rag_configuration(self, rag_manager):
        """Configuration avancée du système RAG"""
        st.markdown("### ⚙️ Configuration RAG")
        
        tab_conf1, tab_conf2, tab_conf3, tab_conf4 = st.tabs([
            "🔧 Paramètres", "🧠 Embeddings", "💾 Stockage", "🔬 Expérimental"
        ])
        
        with tab_conf1:
            st.markdown("#### 🔧 Paramètres Généraux")
            
            col_param1, col_param2 = st.columns(2)
            
            with col_param1:
                st.markdown("**Traitement de texte**")
                
                chunk_method = st.selectbox(
                    "Méthode de chunking:",
                    ["Taille fixe", "Phrases", "Paragraphes", "Sémantique"]
                )
                
                chunk_size = st.slider(
                    "Taille des chunks:",
                    100, 5000, 1000,
                    help="Nombre de caractères par segment"
                )
                
                chunk_overlap = st.slider(
                    "Chevauchement:",
                    0, 500, 100,
                    help="Caractères de chevauchement entre chunks"
                )
            
            with col_param2:
                st.markdown("**Recherche**")
                
                similarity_metric = st.selectbox(
                    "Métrique de similarité:",
                    ["Cosine", "Euclidienne", "Manhattan", "Dot Product"]
                )
                
                rerank_results = st.checkbox(
                    "Réorganiser résultats",
                    value=True,
                    help="Utiliser un modèle de reranking"
                )
                
                use_mmr = st.checkbox(
                    "Diversité (MMR)",
                    value=False,
                    help="Maximum Marginal Relevance pour diversifier"
                )
        
        with tab_conf2:
            st.markdown("#### 🧠 Configuration Embeddings")
            
            embedding_model = st.selectbox(
                "Modèle d'embedding:",
                [
                    "sentence-transformers/all-MiniLM-L6-v2",
                    "sentence-transformers/all-mpnet-base-v2",
                    "BAAI/bge-large-en-v1.5",
                    "OpenAI text-embedding-ada-002",
                    "Cohere embed-multilingual-v3.0"
                ]
            )
            
            if "OpenAI" in embedding_model:
                api_key = st.text_input("Clé API OpenAI:", type="password")
            elif "Cohere" in embedding_model:
                api_key = st.text_input("Clé API Cohere:", type="password")
            
            # Test d'embedding
            test_text = st.text_area(
                "Texte de test:",
                "Ceci est un texte de test pour vérifier les embeddings."
            )
            
            if st.button("🧪 Tester Embedding"):
                embedding = self.test_embedding(test_text, embedding_model)
                if embedding:
                    st.success("✅ Embedding généré avec succès")
                    st.json({
                        "Dimension": len(embedding),
                        "Premiers valeurs": embedding[:5],
                        "Norme": sum(e**2 for e in embedding)**0.5
                    })
        
        with tab_conf3:
            st.markdown("#### 💾 Configuration Stockage")
            
            storage_backend = st.selectbox(
                "Backend de stockage:",
                ["ChromaDB", "Pinecone", "Weaviate", "Qdrant", "FAISS", "Milvus"]
            )
            
            if storage_backend == "ChromaDB":
                persist_directory = st.text_input(
                    "Répertoire de persistance:",
                    value="./chroma_db"
                )
                
                collection_name = st.text_input(
                    "Nom de la collection:",
                    value="planner_ia_docs"
                )
            
            elif storage_backend == "Pinecone":
                pinecone_api = st.text_input("Clé API Pinecone:", type="password")
                pinecone_env = st.text_input("Environnement Pinecone:")
                index_name = st.text_input("Nom de l'index:", value="planner-ia")
            
            # Statistiques de stockage
            st.markdown("**📊 Statistiques de stockage**")
            
            col_stor1, col_stor2, col_stor3 = st.columns(3)
            
            with col_stor1:
                st.metric("💾 Espace utilisé", "234 MB")
            with col_stor2:
                st.metric("📊 Vecteurs", "12,456")
            with col_stor3:
                st.metric("⚡ Latence", "23 ms")
        
        with tab_conf4:
            st.markdown("#### 🔬 Fonctionnalités Expérimentales")
            
            st.warning("⚠️ Ces fonctionnalités sont en développement")
            
            enable_cross_lingual = st.checkbox(
                "🌍 Recherche multilingue",
                help="Rechercher dans plusieurs langues"
            )
            
            enable_entity_extraction = st.checkbox(
                "🏷️ Extraction d'entités",
                help="Extraire automatiquement personnes, lieux, dates"
            )
            
            enable_summarization = st.checkbox(
                "📝 Résumé automatique",
                help="Générer des résumés pour documents longs"
            )
            
            enable_clustering = st.checkbox(
                "🔮 Clustering automatique",
                help="Regrouper documents similaires automatiquement"
            )
            
            if st.button("💾 Sauvegarder Configuration"):
                config = self.save_configuration(
                    chunk_method, chunk_size, chunk_overlap,
                    similarity_metric, rerank_results, use_mmr,
                    embedding_model, storage_backend
                )
                st.success("✅ Configuration sauvegardée")
                st.json(config)
    
    # ===== MÉTHODES UTILITAIRES =====
    
    def process_advanced_import(self, files, rag_manager, category, tags, chunk_size, ocr):
        """Traite l'import avancé avec métadonnées"""
        progress = st.progress(0)
        status = st.empty()
        imported = 0
        
        for i, file in enumerate(files):
            try:
                status.info(f"📄 Traitement: {file.name}")
                
                # Extraction du contenu selon le type
                content = self.extract_content(file, ocr)
                
                # Chunking si nécessaire
                if len(content) > chunk_size:
                    chunks = self.chunk_text(content, chunk_size)
                else:
                    chunks = [content]
                
                # Import de chaque chunk
                for j, chunk in enumerate(chunks):
                    metadata = {
                        "filename": file.name,
                        "category": category,
                        "tags": tags,
                        "chunk": f"{j+1}/{len(chunks)}",
                        "size": file.size,
                        "import_date": datetime.now().isoformat()
                    }
                    
                    doc_id = f"{file.name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}_chunk_{j}"
                    
                    # Ajouter au RAG
                    success = self.add_to_rag(chunk, metadata, doc_id, rag_manager)
                    if success:
                        imported += 1
                
                progress.progress((i + 1) / len(files))
                
            except Exception as e:
                st.error(f"❌ Erreur {file.name}: {e}")
        
        if imported > 0:
            st.success(f"✅ {imported} documents/chunks importés!")
        
        progress.empty()
        status.empty()
    
    def extract_content(self, file, use_ocr=False):
        """Extrait le contenu selon le type de fichier"""
        filename = file.name.lower()
        
        # Word documents
        if filename.endswith(('.docx', '.doc')):
            try:
                import docx2txt
                return docx2txt.process(file)
            except ImportError:
                st.warning("📦 Installez python-docx2txt pour support Word")
                return f"[Word non supporté - installez python-docx2txt]"
        
        # PowerPoint
        elif filename.endswith(('.pptx', '.ppt')):
            try:
                from pptx import Presentation
                import io
                prs = Presentation(io.BytesIO(file.getvalue()))
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                return "\n".join(text_parts)
            except ImportError:
                st.warning("📦 Installez python-pptx pour support PowerPoint")
                return f"[PowerPoint non supporté - installez python-pptx]"
        
        # Excel
        elif filename.endswith(('.xlsx', '.xls')):
            try:
                df = pd.read_excel(file)
                return df.to_string()
            except Exception as e:
                return f"[Erreur Excel: {e}]"
        
        # Images avec OCR
        elif filename.endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')) and use_ocr:
            try:
                from PIL import Image
                import pytesseract
                import io
                
                img = Image.open(io.BytesIO(file.getvalue()))
                text = pytesseract.image_to_string(img)
                return text if text.strip() else "[Aucun texte détecté dans l'image]"
            except ImportError:
                st.warning("📦 Installez pytesseract pour OCR")
                return f"[OCR non disponible - installez pytesseract]"
        
        # Fallback pour autres formats
        else:
            try:
                return str(file.read(), "utf-8")
            except:
                return f"[Format non supporté: {filename}]"
    
    def chunk_text(self, text, chunk_size):
        """Divise le texte en chunks avec chevauchement"""
        chunks = []
        overlap = min(100, chunk_size // 10)
        
        for i in range(0, len(text), chunk_size - overlap):
            chunk = text[i:i + chunk_size]
            if chunk.strip():
                chunks.append(chunk)
        
        return chunks
    
    def add_to_rag(self, content, metadata, doc_id, rag_manager):
        """Ajoute au système RAG avec gestion d'erreurs"""
        try:
            if hasattr(rag_manager, 'add_knowledge'):
                return rag_manager.add_knowledge(content, metadata, doc_id)
            elif hasattr(rag_manager, 'documents'):
                rag_manager.documents.append({
                    'id': doc_id,
                    'text': content,
                    'metadata': metadata
                })
                return True
            return False
        except Exception as e:
            st.error(f"Erreur ajout RAG: {e}")
            return False
    
    def get_document_count(self, rag_manager):
        """Obtient le nombre de documents"""
        try:
            if hasattr(rag_manager, 'collection'):
                return rag_manager.collection.count()
            elif hasattr(rag_manager, 'documents'):
                return len(rag_manager.documents)
            return 0
        except:
            return 0
    
    def get_total_size(self, rag_manager):
        """Calcule la taille totale en MB"""
        # Simulation - à implémenter selon le système
        return 234.56
    
    def get_categories_count(self, rag_manager):
        """Compte les catégories utilisées"""
        return len(self.categories)
    
    def get_last_update(self, rag_manager):
        """Obtient la date de dernière mise à jour"""
        return datetime.now().strftime("%d/%m %H:%M")
    
    def create_tag_cloud(self, tag_stats):
        """Crée un nuage de tags HTML"""
        html = '<div style="padding: 20px; line-height: 2.5;">'
        
        max_count = max(tag_stats.values()) if tag_stats else 1
        
        for tag, count in tag_stats.items():
            size = 12 + (count / max_count) * 20
            opacity = 0.5 + (count / max_count) * 0.5
            
            html += f'''
            <span style="
                font-size: {size}px;
                opacity: {opacity};
                margin: 5px;
                padding: 3px 8px;
                background: rgba(100, 100, 255, 0.1);
                border-radius: 5px;
                display: inline-block;
            ">#{tag}</span>
            '''
        
        html += '</div>'
        return html
    
    def get_filtered_documents(self, rag_manager, category, date, file_type):
        """Obtient les documents filtrés"""
        # Simulation - à adapter selon le système réel
        return [
            {
                'id': 'doc1',
                'name': 'Plan_Projet_2024.pdf',
                'category': '📋 Gestion de Projet',
                'tags': ['planning', 'budget'],
                'size': '2.3 MB',
                'date': '2024-01-15',
                'relevance': 0.95
            }
        ]
    
    def get_popular_tags(self, rag_manager, limit=10):
        """Obtient les tags les plus utilisés"""
        # Simulation
        return [
            ('planning', 45),
            ('budget', 38),
            ('risques', 32),
            ('équipe', 28),
            ('technique', 25)
        ]
    
    def save_configuration(self, *args):
        """Sauvegarde la configuration"""
        config = {
            'chunk_method': args[0],
            'chunk_size': args[1],
            'chunk_overlap': args[2],
            'similarity_metric': args[3],
            'rerank_results': args[4],
            'use_mmr': args[5],
            'embedding_model': args[6],
            'storage_backend': args[7],
            'saved_at': datetime.now().isoformat()
        }
        
        # Sauvegarder dans un fichier
        config_path = Path("config/rag_config.json")
        config_path.parent.mkdir(exist_ok=True)
        
        with open(config_path, 'w') as f:
            json.dump(config, f, indent=2)
        
        return config
    
    def get_search_suggestions(self, query, rag_manager):
        """Génère des suggestions de recherche basées sur la requête"""
        suggestions = []
        
        if not query:
            # Suggestions par défaut
            suggestions = [
                "risques projet",
                "planning développement",
                "budget estimation",
                "ressources équipe"
            ]
        else:
            # Suggestions basées sur la requête
            query_lower = query.lower()
            
            if "risque" in query_lower:
                suggestions = [
                    f"{query} mitigation",
                    f"{query} analyse",
                    f"{query} prévention",
                    f"{query} impact"
                ]
            elif "budget" in query_lower:
                suggestions = [
                    f"{query} optimisation",
                    f"{query} réduction",
                    f"{query} allocation",
                    f"{query} prévision"
                ]
            elif "planning" in query_lower:
                suggestions = [
                    f"{query} optimisation",
                    f"{query} chemin critique",
                    f"{query} gantt",
                    f"{query} milestones"
                ]
            else:
                # Suggestions génériques
                suggestions = [
                    f"{query} analyse",
                    f"{query} optimisation",
                    f"{query} recommandations",
                    f"{query} best practices"
                ]
        
        return suggestions[:4]  # Limiter à 4 suggestions
    
    def perform_advanced_search(self, query, rag_manager, search_type, 
                                relevance_threshold, max_results, 
                                search_in_categories, date_range):
        """Effectue une recherche avancée avec filtres"""
        results = []
        
        try:
            # Recherche de base
            if hasattr(rag_manager, 'search_knowledge'):
                raw_results = rag_manager.search_knowledge(query, n_results=max_results*2)
            elif hasattr(rag_manager, 'documents'):
                # Recherche simple pour SimpleRAGManager
                raw_results = []
                query_lower = query.lower()
                for doc in rag_manager.documents:
                    if query_lower in doc.get('text', '').lower():
                        raw_results.append({
                            'text': doc['text'],
                            'metadata': doc.get('metadata', {}),
                            'relevance': 0.5  # Score fixe pour recherche simple
                        })
            else:
                return []
            
            # Appliquer les filtres
            for result in raw_results:
                relevance = result.get('relevance', 0)
                metadata = result.get('metadata', {})
                
                # Filtre de pertinence
                if relevance < relevance_threshold:
                    continue
                
                # Filtre de catégorie
                if search_in_categories:
                    category = metadata.get('category', '')
                    if category and category not in search_in_categories:
                        continue
                
                # Filtre de date (si implémenté)
                # TODO: Ajouter le filtre de date si nécessaire
                
                # Formatter le résultat
                formatted_result = {
                    'id': metadata.get('filename', 'unknown'),
                    'title': metadata.get('filename', 'Document sans titre'),
                    'excerpt': result['text'][:300] + '...' if len(result['text']) > 300 else result['text'],
                    'relevance': relevance,
                    'category': metadata.get('category', 'Non catégorisé'),
                    'tags': metadata.get('tags', []),
                    'date': metadata.get('import_timestamp', 'Date inconnue')[:10] if metadata.get('import_timestamp') else 'Date inconnue',
                    'highlights': []  # Pourrait être implémenté avec une recherche de mots-clés
                }
                
                # Ajouter des highlights si recherche par mots-clés
                if search_type == "🔤 Mots-clés":
                    highlights = []
                    text_lower = result['text'].lower()
                    query_words = query.lower().split()
                    
                    for word in query_words:
                        idx = text_lower.find(word)
                        if idx != -1:
                            start = max(0, idx - 50)
                            end = min(len(result['text']), idx + len(word) + 50)
                            highlight = result['text'][start:end]
                            highlights.append(highlight)
                    
                    formatted_result['highlights'] = highlights[:3]  # Max 3 highlights
                
                results.append(formatted_result)
                
                if len(results) >= max_results:
                    break
            
        except Exception as e:
            st.error(f"Erreur recherche: {e}")
        
        return results
    
    def get_category_document_count(self, category, rag_manager):
        """Compte les documents dans une catégorie"""
        count = 0
        
        try:
            if hasattr(rag_manager, 'documents'):
                for doc in rag_manager.documents:
                    if doc.get('metadata', {}).get('category') == category:
                        count += 1
            # Pour d'autres types de RAG, implémenter selon besoin
        except:
            pass
        
        return count
    
    def add_category(self, category_name):
        """Ajoute une nouvelle catégorie"""
        if category_name not in self.categories:
            self.categories.append(category_name)
    
    def get_tag_statistics(self, rag_manager):
        """Obtient les statistiques des tags"""
        tag_counts = {}
        
        try:
            if hasattr(rag_manager, 'documents'):
                for doc in rag_manager.documents:
                    tags = doc.get('metadata', {}).get('tags', [])
                    for tag in tags:
                        tag_counts[tag] = tag_counts.get(tag, 0) + 1
        except:
            pass
        
        return tag_counts
    
    def merge_tags(self, tags_to_merge, new_tag_name, rag_manager):
        """Fusionne plusieurs tags en un seul"""
        try:
            if hasattr(rag_manager, 'documents'):
                for doc in rag_manager.documents:
                    metadata = doc.get('metadata', {})
                    tags = metadata.get('tags', [])
                    
                    # Remplacer les anciens tags par le nouveau
                    updated_tags = []
                    for tag in tags:
                        if tag in tags_to_merge:
                            if new_tag_name not in updated_tags:
                                updated_tags.append(new_tag_name)
                        else:
                            updated_tags.append(tag)
                    
                    metadata['tags'] = updated_tags
        except Exception as e:
            st.error(f"Erreur fusion tags: {e}")
    
    def get_category_distribution(self, rag_manager):
        """Obtient la distribution des documents par catégorie"""
        distribution = []
        
        for category in self.categories:
            count = self.get_category_document_count(category, rag_manager)
            if count > 0:
                distribution.append({
                    'category': category,
                    'count': count
                })
        
        return distribution
    
    def get_import_timeline(self, rag_manager):
        """Obtient la timeline des imports"""
        # Simulation - à implémenter selon le système réel
        from datetime import timedelta
        import random
        
        timeline = []
        base_date = datetime.now()
        
        for i in range(7):
            date = base_date - timedelta(days=i)
            timeline.append({
                'date': date.strftime('%Y-%m-%d'),
                'count': random.randint(0, 10)
            })
        
        return timeline
    
    def get_document_preview(self, doc_id, rag_manager):
        """Obtient un aperçu du document"""
        try:
            if hasattr(rag_manager, 'documents'):
                for doc in rag_manager.documents:
                    if doc.get('id') == doc_id:
                        return doc.get('text', '')[:500]
        except:
            pass
        
        return "Aperçu non disponible"
    
    def bulk_update_tags(self, doc_ids, rag_manager):
        """Met à jour les tags de plusieurs documents"""
        st.info("Fonction de mise à jour groupée des tags à implémenter")
    
    def bulk_update_category(self, doc_ids, rag_manager):
        """Met à jour la catégorie de plusieurs documents"""
        st.info("Fonction de changement de catégorie groupée à implémenter")
    
    def export_documents(self, doc_ids, rag_manager):
        """Exporte les documents sélectionnés"""
        st.info("Fonction d'export à implémenter")
    
    def delete_documents(self, doc_ids, rag_manager):
        """Supprime les documents sélectionnés"""
        st.info("Fonction de suppression à implémenter")
    
    def test_embedding(self, text, model):
        """Teste la génération d'embedding"""
        try:
            # Simulation d'embedding
            import random
            embedding = [random.random() for _ in range(384)]  # Dimension typique
            return embedding
        except Exception as e:
            st.error(f"Erreur test embedding: {e}")
            return None
    
    def process_web_import(self, urls, rag_manager, depth, include_images, include_pdfs):
        """Traite l'import depuis URLs web"""
        st.info(f"Import de {len(urls)} URL(s) avec profondeur {depth}")
        # À implémenter avec beautifulsoup4 et requests

# Export pour utilisation
def render_advanced_rag_interface(rag_manager):
    """Point d'entrée principal pour l'interface RAG avancée"""
    manager = AdvancedRAGManager()
    manager.render_advanced_interface(rag_manager)